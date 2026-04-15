import json
import os
from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.util import Pt


def generate_ppt(
    plan_file: str,
    slide_images: list[str],
    output_file: str,
) -> str:
    """
    Generate a PowerPoint presentation from slide images.

    Args:
        plan_file: Path to JSON file containing presentation plan
        slide_images: List of paths to slide images in order
        output_file: Path to output PPTX file

    Returns:
        Status message
    """
    # Load presentation plan
    with open(plan_file, "r", encoding="utf-8") as f:
        plan = json.load(f)

    # Determine slide dimensions based on aspect ratio
    aspect_ratio = plan.get("aspect_ratio", "16:9")
    if aspect_ratio == "16:9":
        slide_width = Inches(13.333)
        slide_height = Inches(7.5)
    elif aspect_ratio == "4:3":
        slide_width = Inches(10)
        slide_height = Inches(7.5)
    else:
        # Default to 16:9
        slide_width = Inches(13.333)
        slide_height = Inches(7.5)

    # Create presentation with specified dimensions
    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    # Get blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Add each slide image
    slides_info = plan.get("slides", [])

    if not slide_images:
        for slide_info in slides_info:
            _add_text_slide(prs, blank_layout, slide_info)
        prs.save(output_file)
        return f"Successfully generated text-based presentation with {len(slides_info)} slides to {output_file}"

    for i, image_path in enumerate(slide_images):
        if not os.path.exists(image_path):
            return f"Error: Slide image not found: {image_path}"

        # Add a blank slide
        slide = prs.slides.add_slide(blank_layout)

        # Load and process image
        with Image.open(image_path) as img:
            # Convert to RGB if necessary (for PNG with transparency)
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")

            # Calculate dimensions to fill slide while maintaining aspect ratio
            img_width, img_height = img.size
            img_aspect = img_width / img_height
            slide_aspect = slide_width / slide_height

            # Convert to EMU for calculations
            slide_width_emu = int(slide_width)
            slide_height_emu = int(slide_height)

            if img_aspect > slide_aspect:
                # Image is wider - fit to width
                new_width_emu = slide_width_emu
                new_height_emu = int(slide_width_emu / img_aspect)
                left = Inches(0)
                top = Inches((slide_height_emu - new_height_emu) / 914400)
            else:
                # Image is taller - fit to height
                new_height_emu = slide_height_emu
                new_width_emu = int(slide_height_emu * img_aspect)
                left = Inches((slide_width_emu - new_width_emu) / 914400)
                top = Inches(0)

            # Save processed image to bytes
            img_bytes = BytesIO()
            img.save(img_bytes, format="JPEG", quality=95)
            img_bytes.seek(0)

            # Add image to slide
            slide.shapes.add_picture(
                img_bytes, left, top, Inches(new_width_emu / 914400), Inches(new_height_emu / 914400)
            )

        # Add speaker notes if available in plan
        if i < len(slides_info):
            slide_info = slides_info[i]
            notes = []

            if slide_info.get("title"):
                notes.append(f"Title: {slide_info['title']}")

            if slide_info.get("subtitle"):
                notes.append(f"Subtitle: {slide_info['subtitle']}")

            if slide_info.get("key_points"):
                notes.append("Key Points:")
                for point in slide_info["key_points"]:
                    notes.append(f"  • {point}")

            if notes:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                if text_frame is not None:
                    text_frame.text = "\n".join(notes)

    # Save presentation
    prs.save(output_file)

    return f"Successfully generated presentation with {len(slide_images)} slides to {output_file}"


def _add_text_slide(prs: Presentation, blank_layout, slide_info: dict) -> None:
    slide = prs.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = slide_info.get("title", "Untitled Slide")
    title_paragraph.font.size = Pt(26)
    title_paragraph.font.bold = True
    title_paragraph.alignment = PP_ALIGN.LEFT

    subtitle = slide_info.get("subtitle")
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.text = subtitle
        subtitle_paragraph.font.size = Pt(16)

    body_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.3), Inches(4.3))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True

    key_points = slide_info.get("key_points") or []
    visual_description = slide_info.get("visual_description")
    content_lines = []
    if key_points:
        content_lines.extend([f"• {point}" for point in key_points])
    if visual_description:
        content_lines.append("")
        content_lines.append(f"Visual direction: {visual_description}")
    if not content_lines:
        content_lines.append("No structured slide content was provided.")

    first = True
    for line in content_lines:
        paragraph = body_frame.paragraphs[0] if first else body_frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(18 if line.startswith("•") else 14)
        first = False


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Generate PowerPoint presentation from slide images"
    )
    parser.add_argument(
        "--plan-file",
        required=True,
        help="Absolute path to JSON presentation plan file",
    )
    parser.add_argument(
        "--slide-images",
        nargs="*",
        required=False,
        default=[],
        help="Absolute paths to slide images in order (space-separated)",
    )
    parser.add_argument(
        "--output-file",
        required=True,
        help="Output path for generated PPTX file",
    )

    args = parser.parse_args()

    try:
        print(
            generate_ppt(
                args.plan_file,
                args.slide_images,
                args.output_file,
            )
        )
    except Exception as e:
        print(f"Error while generating presentation: {e}")
