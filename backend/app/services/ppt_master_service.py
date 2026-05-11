"""PPT Master Integration Service - AI-driven PPT generation using system models."""

from __future__ import annotations

import concurrent.futures
import json
import logging
import threading
import time
import uuid
from collections.abc import Callable
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from deerflow.config import get_app_config
from deerflow.models.factory import create_chat_model

logger = logging.getLogger(__name__)

STYLE_COLORS = {
    "free": {"primary": "#2196F3", "secondary": "#4CAF50", "accent": "#FF9800", "bg": "#FFFFFF", "text": "#2C3E50"},
    "consulting": {"primary": "#005587", "secondary": "#0076A8", "accent": "#F5A623", "bg": "#FFFFFF", "text": "#1A252F"},
    "dark": {"primary": "#1A1A2E", "secondary": "#16213E", "accent": "#0F3460", "bg": "#0A0A0A", "text": "#FFFFFF"},
    "gradient": {"primary": "#667EEA", "secondary": "#764BA2", "accent": "#00D4FF", "bg": "#1A1A2E", "text": "#FFFFFF"},
}


@dataclass(slots=True)
class PPTTaskState:
    task_id: str
    phase: str
    progress: float
    attempts: int
    error: str | None
    is_fallback: bool
    started_at: float
    updated_at: float
    output_path: str | None


@dataclass(slots=True)
class PPTMasterConfig:
    timeout: int = 120
    max_retries: int = 3
    base_delay: float = 2.0
    retry_multiplier: float = 2.0
    max_delay: float = 60.0
    enable_fallback: bool = True
    task_ttl: int = 86400
    task_dir: str = "/tmp/ppt-master-tasks"


class PPTMasterRetryableError(RuntimeError):
    """Retryable PPT Master error."""


class PPTMasterNonRetryableError(RuntimeError):
    """Non-retryable PPT Master error."""


class PPTMasterCancelledError(RuntimeError):
    """Raised when a PPT task is cancelled."""


class PPTMasterTimeoutError(TimeoutError):
    """Raised when the LLM request exceeds timeout."""


def _hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 6:
        r, g, b = tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
        return RGBColor(r, g, b)
    return RGBColor(0, 0, 0)


class PPTMasterService:
    SKILL_DIR = Path(__file__).parent.parent.parent.parent / "skills" / "ppt-master" / "skills" / "ppt-master"
    PROJECTS_DIR = Path("/tmp/ppt-master-projects")
    OUTPUTS_DIR = Path("/tmp/ppt-master-outputs")

    def __init__(self):
        self.projects_dir = self.PROJECTS_DIR
        self.outputs_dir = self.OUTPUTS_DIR
        self.projects_dir.mkdir(parents=True, exist_ok=True)
        self.outputs_dir.mkdir(parents=True, exist_ok=True)
        self._task_lock = threading.RLock()
        self._task_states: dict[str, PPTTaskState] = {}
        self._cancelled_tasks: set[str] = set()
        self._refresh_config()
        self._cleanup_old_tasks()

    def _refresh_config(self) -> None:
        raw_config: dict[str, Any] = {}
        try:
            app_config = get_app_config()
            if app_config.model_extra:
                raw_config = dict(app_config.model_extra.get("ppt_master") or {})
        except Exception:
            logger.warning("Failed to load ppt_master config from config.yaml, using defaults", exc_info=True)

        try:
            self.config = PPTMasterConfig(**raw_config)
        except Exception:
            logger.warning("Invalid ppt_master config detected, using defaults: %s", raw_config, exc_info=True)
            self.config = PPTMasterConfig()

        self.task_dir = Path(self.config.task_dir)
        self.task_dir.mkdir(parents=True, exist_ok=True)

    def _ensure_skill_dir(self) -> bool:
        if not self.SKILL_DIR.exists():
            return False
        required_dirs = ["scripts", "templates", "references"]
        for d in required_dirs:
            if not (self.SKILL_DIR / d).exists():
                return False
        return True

    def _task_file_path(self, task_id: str) -> Path:
        return self.task_dir / f"{task_id}.json"

    def _save_task_state(self, state: PPTTaskState) -> PPTTaskState:
        with self._task_lock:
            state.updated_at = time.time()
            self._task_states[state.task_id] = state
            self._task_file_path(state.task_id).write_text(json.dumps(asdict(state), ensure_ascii=False, indent=2), encoding="utf-8")
            return state

    def _load_task_state(self, task_id: str) -> PPTTaskState | None:
        with self._task_lock:
            state = self._task_states.get(task_id)
            if state is not None:
                if time.time() - state.updated_at > self.config.task_ttl:
                    self._delete_task_state(task_id)
                    return None
                return state

            task_file = self._task_file_path(task_id)
            if not task_file.exists():
                return None

            try:
                state = PPTTaskState(**json.loads(task_file.read_text(encoding="utf-8")))
            except Exception:
                logger.warning("Failed to load PPT task state for %s", task_id, exc_info=True)
                task_file.unlink(missing_ok=True)
                return None

            if time.time() - state.updated_at > self.config.task_ttl:
                self._delete_task_state(task_id)
                return None

            self._task_states[task_id] = state
            return state

    def _delete_task_state(self, task_id: str) -> None:
        self._task_states.pop(task_id, None)
        self._cancelled_tasks.discard(task_id)
        self._task_file_path(task_id).unlink(missing_ok=True)

    def _cleanup_old_tasks(self) -> None:
        self._refresh_config()
        cutoff = time.time() - self.config.task_ttl
        with self._task_lock:
            for task_id, state in list(self._task_states.items()):
                if state.updated_at < cutoff:
                    self._delete_task_state(task_id)

            for task_file in self.task_dir.glob("*.json"):
                try:
                    payload = json.loads(task_file.read_text(encoding="utf-8"))
                    updated_at = float(payload.get("updated_at", 0))
                except Exception:
                    logger.warning("Failed to inspect PPT task file %s during cleanup", task_file, exc_info=True)
                    task_file.unlink(missing_ok=True)
                    continue

                if updated_at < cutoff:
                    task_id = payload.get("task_id") or task_file.stem
                    self._delete_task_state(task_id)

    def _update_task_state(self, task_id: str, **changes: Any) -> PPTTaskState | None:
        state = self._load_task_state(task_id)
        if state is None:
            return None

        for key, value in changes.items():
            setattr(state, key, value)
        return self._save_task_state(state)

    def _ensure_not_cancelled(self, task_id: str | None) -> None:
        if not task_id:
            return
        with self._task_lock:
            if task_id in self._cancelled_tasks:
                raise PPTMasterCancelledError("Task cancelled by user")

    def _classify_error(self, error: Exception) -> str:
        status_code = getattr(error, "status_code", None)
        if status_code is None:
            response = getattr(error, "response", None)
            status_code = getattr(response, "status_code", None)

        if isinstance(status_code, int):
            if status_code == 429 or status_code >= 500:
                return "retryable"
            if status_code in {400, 401, 403, 404}:
                return "non_retryable"

        name = error.__class__.__name__.lower()
        message = str(error).lower()
        haystack = f"{name} {message}"
        retryable_tokens = [
            "readtimeout",
            "timeout",
            "timed out",
            "ratelimit",
            "rate limit",
            "too many requests",
            "service unavailable",
            "bad gateway",
            "gateway timeout",
            "connection reset",
            "temporarily unavailable",
            "server error",
            "internal server error",
        ]
        non_retryable_tokens = [
            "autherror",
            "authentication",
            "unauthorized",
            "permission",
            "forbidden",
            "invalidrequest",
            "invalid request",
            "badrequest",
            "bad request",
            "malformed",
        ]

        if any(token in haystack for token in retryable_tokens):
            return "retryable"
        if any(token in haystack for token in non_retryable_tokens):
            return "non_retryable"
        if isinstance(error, (ConnectionError, TimeoutError, concurrent.futures.TimeoutError)):
            return "retryable"
        return "non_retryable"

    def _sleep_with_cancellation(self, delay: float, task_id: str | None) -> None:
        remaining = delay
        while remaining > 0:
            self._ensure_not_cancelled(task_id)
            interval = min(0.5, remaining)
            time.sleep(interval)
            remaining -= interval

    def _retry_with_backoff(self, operation: Callable[[], Any], *, task_id: str | None = None, operation_name: str = "operation") -> Any:
        last_error: Exception | None = None
        for attempt in range(1, self.config.max_retries + 1):
            self._ensure_not_cancelled(task_id)
            if task_id:
                self._update_task_state(task_id, attempts=attempt)

            try:
                return operation()
            except PPTMasterCancelledError:
                raise
            except Exception as exc:
                last_error = exc
                classification = self._classify_error(exc)
                if task_id:
                    self._update_task_state(task_id, error=str(exc), attempts=attempt)

                if classification != "retryable" or attempt >= self.config.max_retries:
                    raise exc

                delay = min(self.config.base_delay * (self.config.retry_multiplier ** (attempt - 1)), self.config.max_delay)
                logger.warning(
                    "PPT Master %s failed for task %s on attempt %s/%s, retrying in %.1fs: %s",
                    operation_name,
                    task_id,
                    attempt,
                    self.config.max_retries,
                    delay,
                    exc,
                )
                self._sleep_with_cancellation(delay, task_id)

        if last_error is not None:
            raise last_error
        raise PPTMasterRetryableError(f"PPT Master {operation_name} failed without an explicit exception")

    def _extract_json_payload(self, content: str) -> dict[str, Any]:
        for line in content.split("\n"):
            stripped = line.strip()
            if stripped.startswith("{") and stripped.endswith("}"):
                return json.loads(stripped)
            if stripped.startswith("```"):
                continue

        start = content.find("{")
        end = content.rfind("}") + 1
        if start >= 0 and end > start:
            return json.loads(content[start:end])
        raise PPTMasterNonRetryableError("LLM did not return valid JSON outline content")

    def _call_llm_with_retry(self, prompt: str, *, task_id: str | None = None) -> dict[str, Any]:
        def _invoke() -> dict[str, Any]:
            def _call_model() -> Any:
                model = create_chat_model()
                return model.invoke(prompt)

            executor = concurrent.futures.ThreadPoolExecutor(max_workers=1)
            future = executor.submit(_call_model)
            try:
                response = future.result(timeout=self.config.timeout)
            except concurrent.futures.TimeoutError as exc:
                future.cancel()
                executor.shutdown(wait=False, cancel_futures=True)
                logger.warning("PPT Master LLM call timed out after %ss for task %s", self.config.timeout, task_id)
                raise PPTMasterTimeoutError(f"LLM timeout after {self.config.timeout}s") from exc
            except Exception:
                executor.shutdown(wait=False)
                raise
            else:
                executor.shutdown(wait=False)

            content = response.content if hasattr(response, "content") else str(response)
            return self._extract_json_payload(content)

        return self._retry_with_backoff(_invoke, task_id=task_id, operation_name="outline generation")

    def _validate_outline(self, outline: dict[str, Any]) -> bool:
        if not isinstance(outline, dict):
            return False

        if not isinstance(outline.get("title"), str) or not outline.get("title", "").strip():
            return False

        slides = outline.get("slides")
        if not isinstance(slides, list) or len(slides) < 3:
            return False

        for slide in slides:
            if not isinstance(slide, dict):
                return False
            title = slide.get("title")
            if not isinstance(title, str) or not title.strip():
                return False

            slide_type = str(slide.get("type", "content"))
            bullets = slide.get("bullets", [])
            if slide_type != "title":
                if not isinstance(bullets, list) or not any(isinstance(item, str) and item.strip() for item in bullets):
                    return False

        return True

    def _generate_outline(self, topic: str, num_slides: int, style: str, *, task_id: str | None = None) -> tuple[dict[str, Any], bool, str | None]:
        prompt = f"""Create a detailed outline for a {num_slides}-slide presentation on: {topic}

Return a JSON object with this exact structure:
{{
    "title": "Presentation Title",
    "subtitle": "Subtitle or description",
    "slides": [
        {{
            "type": "title|content|conclusion",
            "title": "Slide Title",
            "bullets": ["Point 1", "Point 2", "Point 3"]
        }}
    ]
}}

Style: {style}
Must be exactly {num_slides} slides.
Return ONLY valid JSON, no markdown code blocks."""

        try:
            outline = self._call_llm_with_retry(prompt, task_id=task_id)
            if not self._validate_outline(outline):
                raise PPTMasterNonRetryableError("Generated outline failed validation: slides must be >= 3 and include title/bullets")
            return outline, False, None
        except PPTMasterCancelledError:
            raise
        except Exception as exc:
            logger.warning("PPT Master outline generation failed for task %s, falling back to basic template: %s", task_id, exc, exc_info=True)
            warning = f"{exc}. Used basic template."
            if not self.config.enable_fallback:
                raise
            return self._fallback_outline(topic, num_slides), True, warning

    def _fallback_outline(self, topic: str, num_slides: int) -> dict[str, Any]:
        slides = [{"type": "title", "title": topic, "bullets": []}]
        topics = ["概述", "主要内容", "关键要点", "优势分析", "实施方案"]
        for t in topics[: num_slides - 2]:
            slides.append({"type": "content", "title": t, "bullets": [f"{t}要点1", f"{t}要点2", f"{t}要点3"]})
        slides.append({"type": "conclusion", "title": "总结", "bullets": ["核心结论", "下一步建议", "行动计划"]})
        return {"title": topic, "subtitle": "AI Generated", "slides": slides[:num_slides]}

    def generate(
        self,
        topic: str,
        content: str | None = None,
        source_file: str | None = None,
        num_slides: int = 8,
        style: str = "free",
        template: str | None = None,
        format: str = "ppt169",
        aspect_ratio: str = "16:9",
    ) -> dict[str, Any]:
        del content, source_file, template, format

        self._refresh_config()
        self._cleanup_old_tasks()

        task_id = str(uuid.uuid4())[:8]
        colors = STYLE_COLORS.get(style, STYLE_COLORS["free"])
        output_path = self.outputs_dir / f"presentation_{task_id}.pptx"
        output_path.parent.mkdir(parents=True, exist_ok=True)
        now = time.time()

        self._save_task_state(
            PPTTaskState(
                task_id=task_id,
                phase="outline",
                progress=0.05,
                attempts=0,
                error=None,
                is_fallback=False,
                started_at=now,
                updated_at=now,
                output_path=str(output_path),
            )
        )

        try:
            outline, is_fallback, warning = self._generate_outline(topic, num_slides, style, task_id=task_id)
            self._update_task_state(task_id, phase="building", progress=0.55, is_fallback=is_fallback, error=warning, output_path=str(output_path))

            prs = Presentation()
            prs.slide_width = Inches(13.333) if aspect_ratio == "16:9" else Inches(10)
            prs.slide_height = Inches(7.5)
            blank_layout = prs.slide_layouts[6]

            slides = outline.get("slides", [])
            total_slides = max(len(slides), 1)
            for index, slide_data in enumerate(slides, start=1):
                self._ensure_not_cancelled(task_id)
                slide_type = slide_data.get("type", "content")
                title = slide_data.get("title", "")
                bullets = slide_data.get("bullets", [])

                if slide_type == "title":
                    self._create_title_slide(prs, blank_layout, title, outline.get("subtitle", ""), colors)
                elif slide_type == "conclusion":
                    self._create_conclusion_slide(prs, blank_layout, title, colors)
                else:
                    self._create_content_slide(prs, blank_layout, title, bullets, colors)

                progress = 0.55 + (0.35 * index / total_slides)
                self._update_task_state(task_id, phase="building", progress=progress)

            self._update_task_state(task_id, phase="finalize", progress=0.95)
            self._ensure_not_cancelled(task_id)
            prs.save(str(output_path))

            status = "completed_with_fallback" if is_fallback else "completed"
            message = "PPT generated with basic template" if is_fallback else "PPT generated successfully using AI"
            self._update_task_state(task_id, phase="complete", progress=1.0, error=warning, is_fallback=is_fallback, output_path=str(output_path))

            return {
                "success": True,
                "task_id": task_id,
                "file_path": str(output_path),
                "status": status,
                "message": message,
                "is_fallback": is_fallback,
                "warning": warning,
            }
        except PPTMasterCancelledError as exc:
            logger.warning("PPT Master task %s cancelled", task_id)
            self._update_task_state(task_id, phase="cancelled", error=str(exc))
            return {
                "success": False,
                "task_id": task_id,
                "file_path": None,
                "status": "cancelled",
                "message": "PPT generation cancelled",
                "error": str(exc),
                "is_fallback": False,
                "warning": None,
            }
        except Exception as exc:
            logger.warning("PPT Master task %s failed", task_id, exc_info=True)
            self._update_task_state(task_id, phase="failed", error=str(exc))
            return {
                "success": False,
                "task_id": task_id,
                "file_path": None,
                "status": "failed",
                "message": "PPT generation failed",
                "error": str(exc),
                "is_fallback": False,
                "warning": None,
            }

    def _create_title_slide(self, prs: Presentation, blank_layout, title: str, subtitle: str, colors: dict) -> None:
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = _hex_to_rgb(colors["primary"])
        bg.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = _hex_to_rgb(colors["text"])
        p.alignment = PP_ALIGN.CENTER

        if subtitle:
            sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
            tf = sub.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(28)
            p.font.color.rgb = _hex_to_rgb(colors["text"])
            p.alignment = PP_ALIGN.CENTER

    def _create_content_slide(self, prs: Presentation, blank_layout, title: str, bullets: list, colors: dict) -> None:
        slide = prs.slides.add_slide(blank_layout)
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = _hex_to_rgb(colors["primary"])
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = _hex_to_rgb(colors["text"])

        for i, point in enumerate(bullets):
            bullet_box = slide.shapes.add_textbox(Inches(1), Inches(1.8) + Inches(i * 0.9), Inches(11.333), Inches(0.8))
            tf = bullet_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"• {point}"
            p.font.size = Pt(24)
            p.font.color.rgb = _hex_to_rgb(colors["text"])

    def _create_conclusion_slide(self, prs: Presentation, blank_layout, title: str, colors: dict) -> None:
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = _hex_to_rgb(colors["primary"])
        bg.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = _hex_to_rgb(colors["text"])
        p.alignment = PP_ALIGN.CENTER

    def get_task_state(self, task_id: str) -> dict[str, Any] | None:
        self._cleanup_old_tasks()
        state = self._load_task_state(task_id)
        return asdict(state) if state else None

    def cancel_task(self, task_id: str) -> bool:
        self._cleanup_old_tasks()
        state = self._load_task_state(task_id)
        if state is None:
            return False

        with self._task_lock:
            self._cancelled_tasks.add(task_id)

        if state.phase not in {"complete", "failed", "cancelled"}:
            self._update_task_state(task_id, phase="cancelled", error="Task cancelled by user")
        return True

    def get_output_path(self, task_id: str) -> str:
        pptx_file = self.OUTPUTS_DIR / f"presentation_{task_id}.pptx"
        if pptx_file.exists():
            return str(pptx_file)
        return ""

    def list_templates(self) -> list[dict[str, str]]:
        templates = []
        layouts_dir = self.SKILL_DIR / "templates" / "layouts"
        if layouts_dir.exists():
            for template_dir in layouts_dir.iterdir():
                if template_dir.is_dir():
                    templates.append({"name": template_dir.name, "path": str(template_dir)})
        return templates

    def get_template_info(self, template_name: str) -> dict[str, Any] | None:
        template_path = self.SKILL_DIR / "templates" / "layouts" / template_name
        if not template_path.exists():
            return None
        spec_file = template_path / "design_spec.md"
        if spec_file.exists():
            return {"name": template_name, "path": str(template_path), "spec": spec_file.read_text(encoding="utf-8")}
        return {"name": template_name, "path": str(template_path), "spec": ""}


_ppt_master_service: PPTMasterService | None = None


def get_ppt_master_service() -> PPTMasterService:
    global _ppt_master_service
    if _ppt_master_service is None:
        _ppt_master_service = PPTMasterService()
    return _ppt_master_service
