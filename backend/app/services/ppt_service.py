"""PPT Generation Service - creates natively editable PPTX using python-pptx."""

import uuid
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

COLOR_SCHEMES = {
    "gradient-modern": {
        "primary": "#2196F3",
        "secondary": "#4CAF50",
        "accent": "#FF9800",
        "background": "#FFFFFF",
        "text": "#2C3E50",
        "text_light": "#FFFFFF",
    },
    "dark-premium": {
        "primary": "#1A1A2E",
        "secondary": "#16213E",
        "accent": "#0F3460",
        "background": "#0A0A0A",
        "text": "#FFFFFF",
        "text_light": "#FFFFFF",
    },
    "glassmorphism": {
        "primary": "#667EEA",
        "secondary": "#764BA2",
        "accent": "#00D4FF",
        "background": "#1A1A2E",
        "text": "#FFFFFF",
        "text_light": "#FFFFFF",
    },
    "keynote": {
        "primary": "#000000",
        "secondary": "#1D1D1F",
        "accent": "#0071E3",
        "background": "#FFFFFF",
        "text": "#1A1A1A",
        "text_light": "#FFFFFF",
    },
    "minimal-swiss": {
        "primary": "#000000",
        "secondary": "#FFFFFF",
        "accent": "#FF0000",
        "background": "#FFFFFF",
        "text": "#000000",
        "text_light": "#FFFFFF",
    },
    "consulting": {
        "primary": "#005587",
        "secondary": "#0076A8",
        "accent": "#F5A623",
        "background": "#FFFFFF",
        "text": "#1A252F",
        "text_light": "#FFFFFF",
    },
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 6:
        r, g, b = tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
        return RGBColor(r, g, b)
    return RGBColor(0, 0, 0)


def generate_slide_content(topic: str, num_slides: int, style: str) -> list[dict]:
    slides_content = []

    slides_content.append(
        {
            "type": "title",
            "title": topic,
            "subtitle": "Generated Presentation",
        }
    )

    slides_content.append(
        {
            "type": "content",
            "title": "Overview",
            "bullet_points": [
                "Introduction and background",
                "Key concepts and principles",
                "Main discussion points",
                "Summary and conclusions",
            ],
        }
    )

    content_topics = ["Key Concepts", "Main Features", "Benefits and Advantages", "Implementation Approach", "Results and Impact"]
    for content_topic in content_topics[: num_slides - 3]:
        slides_content.append(
            {
                "type": "content",
                "title": content_topic,
                "bullet_points": [
                    f"Point 1 for {content_topic}",
                    f"Point 2 for {content_topic}",
                    f"Point 3 for {content_topic}",
                ],
            }
        )

    if num_slides > 2:
        slides_content.append(
            {
                "type": "conclusion",
                "title": "Summary",
                "subtitle": "Thank you for your attention",
            }
        )

    return slides_content[:num_slides]


def create_title_slide(prs: Presentation, blank_layout, title: str, subtitle: str, style: str) -> None:
    colors = COLOR_SCHEMES.get(style, COLOR_SCHEMES["gradient-modern"])
    slide = prs.slides.add_slide(blank_layout)

    bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = hex_to_rgb(colors["primary"])
    bg_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(colors["text_light"])
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = subtitle
        sub_p.font.size = Pt(28)
        sub_p.font.color.rgb = hex_to_rgb(colors["text_light"])
        sub_p.alignment = PP_ALIGN.CENTER

    accent_line = slide.shapes.add_shape(1, Inches(5.5), Inches(4.5), Inches(2.333), Inches(0.05))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = hex_to_rgb(colors["accent"])
    accent_line.line.fill.background()


def create_content_slide(prs: Presentation, blank_layout, title: str, bullet_points: list[str], style: str) -> None:
    colors = COLOR_SCHEMES.get(style, COLOR_SCHEMES["gradient-modern"])
    slide = prs.slides.add_slide(blank_layout)

    header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = hex_to_rgb(colors["primary"])
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(colors["text_light"])

    bullet_top = Inches(1.8)
    for i, point in enumerate(bullet_points):
        bullet_box = slide.shapes.add_textbox(Inches(1), bullet_top + Inches(i * 0.9), Inches(11.333), Inches(0.8))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {point}"
        p.font.size = Pt(24)
        p.font.color.rgb = hex_to_rgb(colors["text"])

    accent = slide.shapes.add_shape(1, Inches(0), Inches(7.3), Inches(0.15), Inches(0.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = hex_to_rgb(colors["accent"])
    accent.line.fill.background()


def create_conclusion_slide(prs: Presentation, blank_layout, title: str, subtitle: str, style: str) -> None:
    colors = COLOR_SCHEMES.get(style, COLOR_SCHEMES["gradient-modern"])
    slide = prs.slides.add_slide(blank_layout)

    bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = hex_to_rgb(colors["primary"])
    bg_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(colors["text_light"])
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = subtitle
        sub_p.font.size = Pt(28)
        sub_p.font.color.rgb = hex_to_rgb(colors["text_light"])
        sub_p.alignment = PP_ALIGN.CENTER


def generate_ppt(topic: str, num_slides: int = 8, style: str = "gradient-modern", aspect_ratio: str = "16:9", output_dir: str | None = None) -> str:
    if output_dir is None:
        output_dir = "/tmp/ppt-workspaces"

    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    task_id = str(uuid.uuid4())[:8]
    filename = f"presentation_{task_id}.pptx"
    filepath = output_path / filename

    prs = Presentation()

    if aspect_ratio == "16:9":
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    elif aspect_ratio == "4:3":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    slides_content = generate_slide_content(topic, num_slides, style)

    for slide_data in slides_content:
        slide_type = slide_data.get("type", "content")

        if slide_type == "title":
            create_title_slide(prs, blank_layout, slide_data.get("title", topic), slide_data.get("subtitle", ""), style)
        elif slide_type == "conclusion":
            create_conclusion_slide(prs, blank_layout, slide_data.get("title", "Summary"), slide_data.get("subtitle", ""), style)
        else:
            create_content_slide(prs, blank_layout, slide_data.get("title", "Content"), slide_data.get("bullet_points", []), style)

    prs.save(str(filepath))
    return str(filepath)


class PPTService:
    def __init__(self, workspace_dir: str = "/tmp/ppt-workspaces"):
        self.workspace_dir = Path(workspace_dir)
        self.workspace_dir.mkdir(parents=True, exist_ok=True)

    def generate(self, topic: str, num_slides: int = 8, style: str = "gradient-modern", aspect_ratio: str = "16:9") -> dict:
        try:
            if not topic or not topic.strip():
                return {"success": False, "error": "Topic cannot be empty"}

            if num_slides < 3:
                num_slides = 3
            elif num_slides > 20:
                num_slides = 20

            if style not in COLOR_SCHEMES:
                style = "gradient-modern"

            if aspect_ratio not in ["16:9", "4:3"]:
                aspect_ratio = "16:9"

            file_path = generate_ppt(
                topic=topic.strip(),
                num_slides=num_slides,
                style=style,
                aspect_ratio=aspect_ratio,
                output_dir=str(self.workspace_dir),
            )

            task_id = Path(file_path).stem.replace("presentation_", "")

            return {
                "success": True,
                "task_id": task_id,
                "file_path": file_path,
                "status": "completed",
                "message": "PPT generated successfully",
            }

        except Exception as e:
            return {"success": False, "error": str(e), "status": "failed"}

    def get_output_path(self, task_id: str) -> str:
        filepath = self.workspace_dir / f"presentation_{task_id}.pptx"
        if filepath.exists():
            return str(filepath)
        return ""


_ppt_service: PPTService | None = None


def get_ppt_service() -> PPTService:
    global _ppt_service
    if _ppt_service is None:
        _ppt_service = PPTService()
    return _ppt_service
